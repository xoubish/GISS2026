#!/usr/bin/env python
"""Build a fully NATIVE, editable PowerPoint from the final deck.

Unlike export_pptx.py (screenshots), every text, table, and the loop diagram
here are real PowerPoint objects — edit text, colors, and layout freely.
The two Sisyphus drawings are recolored to deck ink on transparency so they
sit on the paper background; paper figures keep their original pixels.

    final/tools/venv/bin/python final/tools/export_editable_pptx.py

Writes ../../Presentation_editable.pptx (one slide per talking beat; builds
that only accumulate text are collapsed onto one slide). Speaker notes are
carried over from js/scenes.js verbatim.
"""

import math
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
FINAL = os.path.dirname(HERE)
REPO = os.path.dirname(FINAL)
ASSETS = os.path.join(FINAL, 'assets')
BUILD = os.path.join(HERE, 'build')
OUT = os.path.join(REPO, 'Presentation_editable.pptx')

# ----------------------------------------------------------- palette (deck.css)
PAPER = RGBColor(0xEE, 0xF4, 0xE8)
PAPER_HI = RGBColor(0xFF, 0xF8, 0xE8)
INK = RGBColor(0x17, 0x20, 0x16)
INK_SOFT = RGBColor(0x31, 0x42, 0x37)
MUTED = RGBColor(0x65, 0x76, 0x6A)
HUMAN = RGBColor(0xD0, 0x6C, 0x4A)
GREEN = RGBColor(0x3F, 0x7B, 0x5D)
RULE = RGBColor(0x8D, 0x9D, 0x90)      # hairlines

SERIF = 'Palatino'

EMW, EMH = 13.333, 7.5                  # slide inches


# ------------------------------------------------------------------- helpers
def inkify(name, ink=(0x17, 0x20, 0x16)):
    """Drawing on white paper -> deck-ink on transparency. Returns new path."""
    os.makedirs(BUILD, exist_ok=True)
    dst = os.path.join(BUILD, name)
    if os.path.exists(dst):
        return dst
    im = Image.open(os.path.join(ASSETS, name)).convert('L')
    lut = [min(255, max(0, int((240 - l) * 255 / 200))) for l in range(256)]
    alpha = im.point(lut)
    out = Image.new('RGBA', im.size, ink + (255,))
    out.putalpha(alpha)
    out.save(dst)
    return dst


def blank_slide(prs, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def textbox(slide, x, y, w, h=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, runs, size=16, color=INK, italic=False, bold=False, align=None,
         before=0, after=6, line=1.2, tracking=None, first=False):
    """Add one paragraph. `runs` is a string or list of (text, style-dict)."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    if align is not None:
        p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, st in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = st.get('font', SERIF)
        f.size = Pt(st.get('size', size))
        f.bold = st.get('bold', bold)
        f.italic = st.get('italic', italic)
        f.color.rgb = st.get('color', color)
        rPr = r._r.get_or_add_rPr()
        tr = st.get('tracking', tracking)
        if tr:
            rPr.set('spc', str(int(tr * 100)))     # tracking in points
        if st.get('sub'):
            rPr.set('baseline', '-25000')
    return p


def kicker(slide, text, x=0.73, y=0.5, w=9.0, size=11):
    tf = textbox(slide, x, y, w)
    para(tf, text.upper(), size=size, color=MUTED, tracking=2.6,
         after=0, first=True)
    return tf


def scenehead(slide, kick, title, x=0.73, y=0.5, w=9.5, title_size=30):
    kicker(slide, kick, x, y, w)
    tf = textbox(slide, x, y + 0.42, w)
    para(tf, title, size=title_size, color=INK, after=0, line=1.06, first=True)


# lead body text: italic ink, bold runs upright — the deck's `aside lead` voice
def lead(tf, runs, size=16, first=False, after=10, color=INK):
    if isinstance(runs, str):
        runs = [(runs, {})]
    styled = []
    for text, st in runs:
        s = {'italic': not st.get('bold') and not st.get('upright'),
             'size': size, 'color': color}
        s.update(st)
        styled.append((text, s))
    return para(tf, styled, size=size, first=first, after=after, line=1.32)


def caption(slide, text, x, y, w, align=None):
    tf = textbox(slide, x, y, w)
    para(tf, text, size=10.5, color=MUTED, italic=True, after=0, line=1.3,
         align=align, first=True)


def picture(slide, path, x, y, w=None, h=None, border=False):
    iw, ih = Image.open(path).size
    if w and not h:
        h = w * ih / iw
    if h and not w:
        w = h * iw / ih
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = RULE
        pic.line.width = Pt(0.75)
    return pic, w, h


# --------------------------------------------------------------- table helpers
NO_STYLE = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # "No Style, No Grid"


def plain_table(slide, x, y, w, col_ws, n_rows):
    gf = slide.shapes.add_table(n_rows, len(col_ws), Inches(x), Inches(y),
                                Inches(w), Inches(0.4 * n_rows))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tblPr = tbl._tbl.tblPr
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = tblPr.makeelement(qn('a:tableStyleId'), {})
        tblPr.append(sid)
    sid.text = NO_STYLE
    total = sum(col_ws)
    for i, cw in enumerate(col_ws):
        tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    for row in tbl.rows:
        for cell in row.cells:
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
    return tbl


def rule_below(cell, color='465B3E', w_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    old = tcPr.find(qn('a:lnB'))
    if old is not None:
        tcPr.remove(old)
    lnB = tcPr.makeelement(qn('a:lnB'), {'w': str(int(w_pt * 12700)),
                                         'cap': 'flat'})
    fill = lnB.makeelement(qn('a:solidFill'), {})
    clr = lnB.makeelement(qn('a:srgbClr'), {'val': color})
    fill.append(clr)
    lnB.append(fill)
    tcPr.insert(0, lnB)


def cell_para(cell, runs, size=12, color=INK, italic=False, tracking=None,
              first=True, after=4):
    tf = cell.text_frame
    tf.word_wrap = True
    return para(tf, runs, size=size, color=color, italic=italic,
                tracking=tracking, first=first, after=after, line=1.25)


def mini_table(slide, x, y, w, heads, cols, cell_size=12.5, head_size=10.5):
    """The deck's mini-table idiom: muted small-cap heads over a hairline,
    italic entries below."""
    tbl = plain_table(slide, x, y, w, [1] * len(heads), 2)
    tbl.rows[0].height = Inches(0.34)
    for i, htext in enumerate(heads):
        c = tbl.cell(0, i)
        cell_para(c, htext.upper(), size=head_size, color=MUTED, tracking=2.2)
        rule_below(c)
    for i, lines in enumerate(cols):
        c = tbl.cell(1, i)
        for j, ln in enumerate(lines):
            cell_para(c, ln, size=cell_size, color=INK_SOFT, italic=True,
                      first=(j == 0), after=6)
    return tbl


# ------------------------------------------------------------- loop diagram
def arrow(slide, x1, y1, x2, y2, dashed=False, color=INK_SOFT, w=1.4):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    if dashed:
        dash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(dash)
    tail = ln.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def loop_diagram(slide, cx=6.9, cy=3.0, rx=4.35, ry=1.72):
    """The handwritten cycle as native shapes: italic lowercase nodes, the
    dashed return from `good enough?` back to `question`."""
    labels = ['question', 'what matters?', 'the ruler', 'data',
              'compression', 'inference', 'answer', 'good enough?']
    angles = [90, 45, 0, -45, -90, -135, 180, 135]
    pos = []
    for a in angles:
        t = math.radians(a)
        pos.append((cx + rx * math.cos(t), cy - ry * math.sin(t)))
    bw, bh = 1.9, 0.38
    for (px, py), lab in zip(pos, labels):
        tf = textbox(slide, px - bw / 2, py - bh / 2, bw, bh)
        tf.word_wrap = False
        para(tf, lab, size=15, color=INK, italic=True, after=0, first=True,
             align=None)
        tf.paragraphs[0].alignment = 2   # PP_ALIGN.CENTER == 2
    for i in range(len(pos)):
        j = (i + 1) % len(pos)
        x1, y1 = pos[i]
        x2, y2 = pos[j]
        dx, dy = x2 - x1, y2 - y1
        last = (j == 0)
        s, e = (0.30, 0.70) if not last else (0.24, 0.76)
        arrow(slide, x1 + dx * s, y1 + dy * s, x1 + dx * e, y1 + dy * e,
              dashed=last)


# =============================================================== build slides
prs = Presentation()
prs.slide_width = Inches(EMW)
prs.slide_height = Inches(EMH)

opening_png = inkify('opening.png')
closing_png = inkify('closing.png')

# ------------------------------------------------------------- 1 · OPENING
s = blank_slide(prs, notes=(
    'GIPS framing · AI/philosophy line · promise: Sisyphus, and all I do is '
    'his story.\n\nScene note: nothing written on screen beyond the title '
    'block; all talking points spoken.'))
tf = textbox(s, 0.73, 0.95, 6.2)
para(tf, 'GIPS 2026', size=12, color=MUTED, tracking=3.4, after=24, first=True)
para(tf, 'Sisyphus,\nOptimizing in a\nDynamic Universe', size=43, color=INK,
     line=1.05, after=28)
para(tf, 'Shoubaneh Hemmati (Caltech/IPAC)', size=13, color=INK_SOFT,
     italic=True, after=0)
picture(s, opening_png, 6.55, 0.28, h=6.95)

# ---------------------------------------------------------------- 2 · MYTH
s = blank_slide(prs, notes=(
    'He climbs while you retell the myth — about fifteen seconds of push. '
    'If you talk longer he holds near the crest; he never reaches it. '
    'The fall lands next: let it settle before speaking again.\n\n'
    'The bridge to the whole talk: pushing toward an optimum is an '
    'optimization. Gradient ascent, if you may.\n\n'
    'Scene note: in the browser deck the pen Sisyphus climbs in real time '
    'and the boulder settles in the bowl below.'))
tf = textbox(s, 0.73, 0.55, 6.0)
para(tf, 'The Myth', size=30, color=INK, after=0, first=True)
tf = textbox(s, 0.73, 1.9, 5.6)
lead(tf, 'Condemned by the gods to push a boulder up a hill for eternity.',
     size=18, first=True, after=18)
lead(tf, '— only for it to roll back, each time, as he nears the top.',
     size=18, after=18)
lead(tf, [('Gradient ascent, or descent with a flipped hill: ', {}),
          ('Sisyphus is solving an optimization problem.', {'bold': True})],
     size=18)
picture(s, opening_png, 7.6, 0.9, h=5.7)

# -------------------------------------------------------- 3 · WHY SISYPHUS
s = blank_slide(prs, notes=(
    'He sits on the summit he never reached, cup in hand. Speak the claim '
    'plainly: almost every decision has his shape — a guess at the world, '
    'then a push toward better.\n\n'
    'The seed of the big table: start with the drink because it is in the '
    'drawing, then the person. The same pair runs in a brain, in ML and '
    'astronomy, and in science as a community.\n\n'
    'Inference: maintaining a model of reality from incomplete, noisy '
    'observations — the guess-at-the-world half.\n\n'
    'Optimization is the mechanical half: align the model with the data, '
    'descend the error surface — the push toward better.'))
scenehead(s, 'Inference & optimization', 'Why We Are Sisyphus')
tf = textbox(s, 0.73, 1.95, 8.2)
lead(tf, 'Almost every decision we make has his shape: a guess at the world, '
     'then a push toward better.', size=17, first=True)
mini_table(
    s, 0.73, 2.9, 10.4,
    ['a brain', 'an AI / model', 'science / astro'],
    [['Should I have another drink?', 'Do I like this person?'],
     ['Is this a cat or a dog?', 'What is the redshift of this galaxy?'],
     ['Reionization — by AGNs or stars?',
      'General Relativity, or modified gravity?']])
tf = textbox(s, 0.73, 5.35, 11.4)
lead(tf, [('Inference', {'bold': True}),
          (' — maintaining a model of reality from incomplete, noisy '
           'observations: p(world | data).', {})], size=15, first=True,
     after=12)
lead(tf, [('Optimization', {'bold': True}),
          (' — the mechanical part: align the model with the data. Descend '
           'the error surface.', {})], size=15)

# ---------------------------------------------------------------- 4 · LOOP
s = blank_slide(prs, notes=(
    'The full loop is already on screen. Walk it once: question, what '
    'matters, ruler, data, compression, inference, answer, good enough?\n\n'
    'Branch one. Not good enough? The absurd answer, 42, means the pipeline '
    'failed somewhere; debug the ruler, the data, the compression, the '
    'model.\n\n'
    'Branch two. Good enough? Stop, or ask the next question. This is the '
    'whole talk: failure sends you inside the pipeline; success ends '
    'locally or restarts the loop.\n\n'
    'Scene note: everything after this scene is a zoom into one of these '
    'stations.'))
kicker(s, 'The loop')
loop_diagram(s)
tf = textbox(s, 1.9, 5.35, 10.2)
lead(tf, [('Not good enough?', {'bold': True}), ('  42 …', {})],
     size=15, first=True, after=2)
p = lead(tf, 'debug the pipeline: the ruler, the data, the compression, …',
         size=14, color=INK_SOFT, after=10)
p.level = 1
lead(tf, [('Good enough?', {'bold': True})], size=15, after=2)
p = lead(tf, 'stop — or ask the next question', size=14, color=INK_SOFT,
         after=0)
p.level = 1

# ------------------------------------------------------- 5 · THREE SYSTEMS
s = blank_slide(prs, notes=(
    'The whole table at once — name the three columns, promise the walk '
    'down the rows.\n\n'
    'The loss: biology, culture and loss aversion for the brain; a '
    'handcrafted mathematical loss for the model; community consensus and '
    'Occam’s razor for science.\n\n'
    'The ruler — biased and bandwidth-limited everywhere. Senses and '
    'attention extended by instruments; sensors, pixel grids, tokenizers; '
    'telescopes, bandpasses, spectrographs. Gamma rays become visible to '
    'the brain through tools, screens and conventions.\n\n'
    'The data: a serial, irreversible sensory stream; fixed, re-shuffled '
    'batches in memory; archived surveys of a non-re-runnable universe.\n\n'
    'The compression: synaptic weights and cortical latent spaces; latent '
    'embeddings and bottlenecks; physical laws and cosmological '
    'parameters.\n\n'
    'The learning signal: prediction error — dopaminergic RPE, surprise; '
    'analytical gradients via backprop; empirical anomalies and measurement '
    'tension.\n\n'
    'Can it restart? The brain cannot — synaptic history is irreversible. '
    'The model can — re-seed, wipe, retrain. Science partially — paradigm '
    'shifts are slow and expensive.\n\n'
    'Who audits the loop? Evolution; we do, with all our proxy blind '
    'spots; peer review and replication — slow, and noisy.\n\n'
    'Pull back: the implementations look different; the loop is '
    'recognizably the same. 12-minute cut: walk rows 1, 5 and 7 only.'))
kicker(s, 'Three systems, one loop')
ROWS = [
    ('The loss', 'Biology, culture & loss aversion',
     'Handcrafted mathematical loss', 'Community consensus & Occam’s razor'),
    ('The ruler', 'Senses, attention & instruments',
     'Sensors, pixel grids & tokenizers',
     'Telescopes, bandpasses, spectrographs…'),
    ('The data', 'Serial, irreversible sensory stream',
     'Fixed, re-shuffled batches in memory',
     'Archived surveys; a non-re-runnable universe'),
    ('The compression', 'Synaptic weights & cortical latent spaces',
     'Latent embeddings & bottlenecks',
     'Physical laws & cosmological parameters'),
    ('The learning signal', 'Prediction error (dopaminergic RPE, surprise)',
     'Analytical gradients via backpropagation',
     'Empirical anomalies & measurement tension'),
    ('Can it restart?', 'No — irreversible synaptic history',
     'Yes — re-seed, wipe weights, retrain',
     'Partially — paradigm shifts are slow & expensive'),
    ('Who audits the loop?', 'Evolution; hardwired biology',
     'We do — with all our proxy blind spots',
     'Peer review & replication (slow and noisy)'),
]
tbl = plain_table(s, 0.73, 1.15, 11.9, [1.55, 2.6, 2.6, 2.6], 8)
tbl.rows[0].height = Inches(0.42)
for r in range(1, 8):
    tbl.rows[r].height = Inches(0.76)
for i, h in enumerate(['dimension', 'brain', 'model / AI',
                       'science as community']):
    c = tbl.cell(0, i)
    cell_para(c, h.upper(), size=11, color=MUTED, tracking=2.0)
    rule_below(c, w_pt=1.0)
for r, row in enumerate(ROWS, start=1):
    c = tbl.cell(r, 0)
    cell_para(c, row[0], size=13, color=INK)
    for i in (1, 2, 3):
        cell_para(tbl.cell(r, i), row[i], size=12.5, color=INK_SOFT,
                  italic=True)
    if r < 7:
        for i in range(4):
            rule_below(tbl.cell(r, i), w_pt=0.5, color='B9C4B4')

# ---------------------------------------------------------------- 6 · INFO
INFO_KICK = 'Information theory · Statistics of the loop'


def info_slide(head_runs, body_paras, notes, first=False):
    s = blank_slide(prs, notes=notes)
    if first:
        scenehead(s, 'Information theory', 'Statistics of the Loop')
        y = 2.2
    else:
        kicker(s, INFO_KICK)
        y = 1.35
    tf = textbox(s, 0.73, y, 8.4)
    lead(tf, head_runs, size=19, first=True, after=14)
    for runs in body_paras:
        lead(tf, runs, size=16, after=12)
    return s


info_slide(
    [('Shannon’s question: how unpredictable is the next observation?',
      {'bold': True})],
    [[('X is the next observation. p(X) is what we expect to see. Entropy, '
       'H(X), is average surprise:', {})],
     [('H(X) = E[−log p(X)] = −Σ p(x) log p(x)', {'upright': True,
                                                  'color': INK_SOFT})],
     [('Low H: little room. High H: room to learn. Usefulness depends '
       'on θ.', {})]],
    notes=(
        'Start with Shannon’s practical question: if observations arrive as '
        'a stream, how unpredictable is the next one? X is the next '
        'observation; p(X) is what your current model expects to see; '
        'entropy H(X) is average surprise. Same family as stat mech '
        'entropy, but this is the message/observation version. Low entropy '
        'means a dimension is almost constant, so it cannot carry much. '
        'High entropy means more possible variation, but that variation '
        'could still be noise. Entropy tells capacity: low H gives little '
        'room; high H gives room to learn. Usefulness depends on θ.\n\n'
        'Scene note (12-minute cut): narrate entropy, Fisher and surprise; '
        'step through the two compression beats and the update quickly; '
        'give the final two beats one sentence each.'),
    first=True)

info_slide(
    [('The decision question — mutual information, I(X;θ).', {'bold': True})],
    [[('How much of X is about θ, the thing I care about?', {})],
     [('Entropy is capacity. Mutual information is relevance.', {})]],
    notes=(
        'The bridge statistic: mutual information. If entropy asks how much '
        'a data dimension can vary, mutual information asks how much of '
        'that variation is about the thing I care about. I(X;θ) is '
        'relevance: how much observing X reduces uncertainty about θ. A '
        'high-entropy channel can be pure noise; a quieter channel can be '
        'decisive if it tracks θ. Worth saying if the room is technical: '
        'this is the SAME quantity as the experimental-design criterion at '
        'the end of the scene — expected information gain about θ from an '
        'observation IS I(θ;Y). Also: “capacity” is loose — in Shannon’s '
        'vocabulary channel capacity is a maximised mutual information, not '
        'an entropy. Fall back on: entropy is how much the reading can '
        'vary, mutual information is how much of that variation is about '
        'θ.'))

info_slide(
    [('The ruler — Fisher information, I(θ).', {'bold': True}),
     (' How much of what I measure is about θ: the same basin, flat under '
      'one ruler, steep under another. Where Fisher is zero, nothing '
      'downstream can ever recover it.', {})],
    [[('Its test: the Cramér–Rao bound, 1/I(θ) — the best any estimator can '
       'ever do through this instrument.', {})]],
    notes=(
        'THE RULER. Fisher is about θ, not about volume: the same basin is '
        'flat under one ruler and steep under another. And the hard '
        'sentence: where Fisher is zero — a wavelength, a conversational '
        'cue — no network and no brain downstream can ever recover it. If '
        'pressed, be precise about which zero is meant: the claim is '
        'airtight for a direction the likelihood does not depend on at all '
        '(the band you never observed), which is the case the examples '
        'are. It is not a claim about Fisher vanishing at one point of an '
        'otherwise informative curve, and the Cramér–Rao bound in its '
        'simple form is a statement about unbiased estimators.'))

info_slide(
    [('The compression — sampling.', {'bold': True}),
     (' Not Shannon’s compression: this one is which numbers you keep. '
      'Five, evenly spaced — and the basin is simply not in them.', {})],
    [[('Nothing was wrong with the data.', {})]],
    notes=(
        'THE COMPRESSION, first lesson. Five numbers, evenly spaced — a '
        'perfectly reasonable summary, and the basin is simply not in it. '
        'The ring is where this summary thinks the minimum is. Nothing was '
        'wrong with the data.'))

info_slide(
    [('The same budget, placed where the curvature lives', {'bold': True}),
     (' — and the basin comes back. Sampling is a ruler too.', {})],
    [[('Its test: does T(D) keep what D knew about θ — I(T(D); θ) against '
       'I(D; θ)?', {})]],
    notes=(
        'Same budget, placed where the curvature lives — and the basin '
        'comes back. No new observation. Sampling is a ruler too: what you '
        'keep decides what exists.'))

info_slide(
    [('The learning signal — surprise, −log p(x).', {'bold': True}),
     (' The datum nothing predicted. Prediction error, gradient, anomaly: '
      'the same engine in all three columns.', {})],
    [[('Its test: converged, yet still surprised — residuals with structure '
       'the fit cannot explain. The optimizer is done, and wrong.', {})]],
    notes=(
        'THE LEARNING SIGNAL. One datum, far from anything the model '
        'predicts: −log p(x). Prediction error in a brain, the gradient in '
        'a model, the anomaly in a field — the same engine, across all '
        'three columns of the table.'))

info_slide(
    [('The update — information gain, D', {'bold': True}),
     ('KL', {'bold': True, 'sub': True}),
     ('.', {'bold': True}),
     (' How much did the bracket close?', {})],
    [[('Surprise is not gain — and when each new datum barely moves the '
       'posterior, this data, under this compression, is spent.', {})]],
    notes=(
        'THE UPDATE. Information gain is the bracket closing: '
        'D_KL(posterior ‖ prior). Two cautions, out loud: low entropy is '
        'not truth, and surprise is not gain — a bizarre datum can teach '
        'you nothing about θ.'))

info_slide(
    [('Finding the bottleneck.', {'bold': True}),
     (' Three checks, cheapest first:', {})],
    [[('Converged but still surprised?', {'bold': True}),
      (' The loss or the model. Change the ruler.', {})],
     [('Raw data beats your summary?', {'bold': True}),
      (' The compression. Recompress — no new observation needed.', {})],
     [('The Cramér–Rao bound itself too wide?', {'bold': True}),
      (' The data. Only then go buy more.', {})]],
    notes=(
        'FINDING THE BOTTLENECK — the three trails out of the basin. Check '
        'in order of cost. Still surprised after convergence → the loss or '
        'model: exchange the landscape. Raw beats the summary → the '
        'compression: re-read the same data. The Cramér–Rao bound itself '
        'too wide → the data: over the crest for more. Say the sting: the '
        'reflex is always “more data” — it is the third check, not the '
        'first.'))

info_slide(
    [('The next observation — expected information gain.', {'bold': True}),
     (' Rank the observations you could take by what you expect to learn. '
      'Buy that one.', {})],
    [[('A little resolution can beat a great many photons.', {})]],
    notes=(
        'THE NEXT QUESTION. And if it truly is the data — experimental '
        'design: which observation do I expect to teach me the most? Rank '
        'them, buy that one. A little resolution can beat a great many '
        'photons. If asked: EIG(d) = E_y D_KL[p(θ|y,d) ‖ p(θ)]. Do not '
        'oversell the ranking — MacKay 1992 (Neural Computation 4:590) is '
        'mostly about how it goes wrong. (1) There is no single “most '
        'informative”: three criteria, depending on whether you want '
        'information about all the parameters, about a named region of '
        'interest, or about telling two models apart — choosing the '
        'criterion is another ruler choice. (2) The naive criterion '
        'reduces to “measure where your error bars are widest”, which '
        'sends you off to the edges of the input space; name the region of '
        'interest first. (3) MacKay’s own Achilles heel: all of it '
        '“estimates the utility of a data point assuming that the model is '
        'correct”, and he closes with “the search for ideal measures of '
        'data utility is still open” — the bottleneck beat again.'))

# ---------------------------------------------------- 7 · EXAMPLE — PERSON
PERSON_KICK = 'Example one · A person you just met'


def person_slide(head_runs, body_paras, notes, first=False):
    s = blank_slide(prs, notes=notes)
    if first:
        scenehead(s, 'Example one', 'A person you just met')
        y = 2.2
    else:
        kicker(s, PERSON_KICK)
        y = 1.35
    tf = textbox(s, 0.73, y, 8.4)
    lead(tf, head_runs, size=19, first=True, after=14)
    for runs in body_paras:
        lead(tf, runs, size=16, after=12)
    return s


person_slide(
    [('Do I like this new person?', {'bold': True}),
     (' Nobody walks in blank: you arrive with a model — a prior built from '
      'every person before them.', {})],
    [],
    notes=(
        'The same ground, retold as a person. The question: do I like '
        'them? Gesture at the resting boulder — the model you arrive with. '
        'Nobody walks in blank: a prior built from every person before '
        'them, context, and yes, stereotype — that is what a prior is.\n\n'
        'Scene note: the audience has already seen every one of these '
        'drawings with statistics captions in scene 6 — now the same '
        'pictures get human captions. That rhyme is the thesis. 12-minute '
        'cut: beats 2, 3, 6, 7.'),
    first=True)

person_slide(
    [('The data.', {'bold': True}),
     (' Words, tone, timing, what they laugh at — an enormous stream, and '
      'most of its entropy is noise.', {})],
    [],
    notes=(
        'The stream: words, tone, timing, what they laugh at, how they '
        'treat the waiter. Enormous entropy — and most of it noise.'))

person_slide(
    [('Distance is resolution.', {'bold': True}),
     (' At first you sample coarsely: job, manners, small talk.', {})],
    [[('“They seem nice” — one blurred pixel.', {})]],
    notes=(
        'Distance is resolution. Across a dinner table you sample coarsely '
        '— job, manners, small talk. “They seem nice” is one blurred '
        'pixel.'))

person_slide(
    [('Closer, the picture resolves.', {'bold': True}),
     (' Finer sampling separates kindness from politeness — two sources '
      'that were one blur.', {})],
    [[('Liking changes with distance because the data does — and the fit '
       'never stops: every encounter is another step of the optimization.',
       {})]],
    notes=(
        'Move closer and the sampling gets finer: shared work, a hard '
        'week, a long trip. Kindness separates from politeness — two '
        'sources that were one blur. And the optimization is constant: '
        'every encounter re-fits the model toward the data. Liking changes '
        'with distance because the data does.'))

person_slide(
    [('Some data carry nothing.', {'bold': True}),
     (' An hour of small talk is pleasant — and Fisher-blind to what you '
      'care about.', {})],
    [[('No amount of it will measure reliability.', {})]],
    notes=(
        'Zero-Fisher data: an hour of small talk is pleasant and carries '
        'nothing about the trait you care about. The flat parabola. No '
        'volume of it will measure reliability.'))

person_slide(
    [('And some data are the same datum again.', {'bold': True}),
     (' The tenth coffee repeats the ninth: plenty of observations, no new '
      'information.', {})],
    [[('The bracket stops closing.', {})]],
    notes=(
        'Duplicate data: the tenth coffee repeats the ninth. Plenty of '
        'observations, no new information — the bracket stops closing. '
        'More data is not more information.'))

person_slide(
    [('Then — surprise.', {'bold': True}),
     (' Something the model never predicted: −log p spikes, and the picture '
      'reorganizes.', {})],
    [[('Sometimes it wasn’t noise. The person moved.', {})]],
    notes=(
        'Then surprise: something the model never predicted. −log p spikes '
        'and the picture reorganizes. And sometimes it was not noise — the '
        'person moved. People are not stationary; the landscape shifts '
        'while you climb it.'))

person_slide(
    [('Good enough?', {'bold': True}),
     (' For “another coffee?” — converged long ago. For “trust them with '
      'what matters” — a different loss entirely, and no number of coffees '
      'helps.', {})],
    [[('It needs a different observation: which encounter teaches the most '
       'is experimental design, about a person.', {})]],
    notes=(
        'Good enough? For “another coffee?” — converged long ago. For '
        '“trust them with what matters” — that is a different loss, and '
        'the same person ranks differently under it; no number of coffees '
        'helps; it needs a different observation entirely: responsibility, '
        'disagreement, stress. Choosing that encounter is experimental '
        'design, about a person — and ethics bound which experiments you '
        'may run.'))

# ----------------------------------------------------- 8 · EXAMPLE — JAISP
JAISP_KICK = 'Example two · JAISP'

s = blank_slide(prs, notes=(
    'The bridge, in one breath: the person resolved because different data '
    'saw them differently. Point it at the sky — Rubin, Euclid, WISE, JWST '
    'are each a lossy projection of the same reality, at their own '
    'distance and resolution.\n\n'
    'Scene note (12-minute cut): receipts get twenty seconds together; '
    'beats 3–5 are the spine, keep all three.'))
scenehead(s, 'Example two — JAISP', 'One foundation, many rulers')
tf = textbox(s, 0.73, 2.3, 8.6)
lead(tf, [('The same lesson, pointed at the sky.', {'bold': True}),
          (' Rubin, Euclid, WISE, JWST — every instrument is a lossy '
           'projection of one sky, at its own distance and resolution.',
           {})], size=19, first=True)

s = blank_slide(prs, notes=(
    'The receipts — say one line each, gesture, move on. Everetts+: Euclid '
    'NISP sharpened with what JWST/NIRCam taught, 5× finer sampling. '
    'Rezaee+: WISE with Spitzer, 4.6×. Haghjoo+: JWST prism spectra to '
    'grating resolution, R 100 → 1000. Each one: one teacher, one student, '
    'one loss, one question.'))
kicker(s, JAISP_KICK)
tf = textbox(s, 0.73, 1.35, 10.5)
lead(tf, [('We did it pairwise first.', {'bold': True}),
          (' One teacher, one student, one loss per question.', {})],
     size=19, first=True)
caps = ['Everetts, Hemmati, et al. — NISP → NIRCam, 5× finer.',
        'Rezaee, Hemmati, et al. — WISE → Spitzer, 4.6× finer.',
        'Haghjoo, Hemmati, et al. — prism → grating, R 100 → 1000.']
for i, (img, cap) in enumerate(zip(
        ['paper_nisp.png', 'paper_wise.png', 'paper_spectra.png'], caps)):
    x = 0.73 + i * 4.15
    _, w, h = picture(s, os.path.join(ASSETS, img), x, 2.5, w=3.8,
                      border=True)
    caption(s, cap, x, 2.5 + h + 0.12, 3.8)

s = blank_slide(prs, notes=(
    'The honest problem: it works, and it does not scale. Every new '
    'question starts again from raw pixels — the compression cost '
    'multiplies with the questions. The general answer: learn the '
    'compression once, before you know the question — a foundation — and '
    'make every task a small head with its own loss.'))
kicker(s, JAISP_KICK)
tf = textbox(s, 0.73, 1.6, 9.2)
lead(tf, [('It works — and it does not scale.', {'bold': True}),
          (' Every new question recompresses the sky from scratch.', {})],
     size=19, first=True, after=14)
lead(tf, [('So learn the compression ', {}), ('once', {'bold': True}),
          (': a foundation — and every task a small head with its own '
           'loss.', {})], size=17)

s = blank_slide(prs, notes=(
    'JAISP. Ten bands, Rubin and Euclid together, one shared latent — '
    'self-supervised, each band predicted from the other nine, each '
    'instrument at its delivered sampling. About nine million parameters. '
    'Detection on held-out sky: 93% complete, 94% pure against the '
    'published VIS catalogue — and 0.45 mag deeper than one band '
    'supports. Say the numbers, not the architecture; the picture is '
    'there so nobody takes them on faith.'))
kicker(s, JAISP_KICK)
tf = textbox(s, 0.73, 1.6, 4.6)
lead(tf, [('JAISP.', {'bold': True}),
          (' Ten bands, one shared latent — self-supervised, '
           '≈9M parameters.', {})], size=17, first=True, after=14)
lead(tf, [('Detection on held-out sky: ', {}),
          ('93% complete, 94% pure', {'bold': True}),
          (' — and ', {}), ('0.45 mag deeper', {'bold': True}),
          (' than one band supports.', {})], size=17)
_, w, h = picture(s, os.path.join(ASSETS, 'jaisp_architecture_crop.png'),
                  5.7, 1.5, w=7.0, border=True)
caption(s, 'Ten bands → two-stream stems → one shared latent → detection · '
        'astrometry · photometry · shape · redshift — each head its own '
        'loss.', 5.7, 1.5 + h + 0.14, 7.0)

s = blank_slide(prs, notes=(
    'The proof, on the θ axis the talk has been walking — drawn to scale, '
    '1 unit = 2 mas. Raw cross-survey scatter, Rubin against VIS: about 50 '
    'mas. A position head reading the frozen latent: 14–17 mas; injected '
    'sources recovered to 19 mas at S/N = 5 — near the floor the VIS '
    'labels themselves set. The latent carried VIS sharpness to '
    'everything tied to it.'))
kicker(s, JAISP_KICK)
tf = textbox(s, 0.73, 1.6, 4.6)
lead(tf, [('One head, its own loss: astrometry.', {'bold': True}),
          (' Raw cross-survey scatter ≈ 50 mas. The head, reading the '
           'frozen latent: ', {}), ('14–17 mas', {'bold': True}),
          ('.', {})], size=17, first=True, after=14)
lead(tf, 'The latent carried VIS sharpness to everything tied to it.',
     size=17)
_, w, h = picture(s, os.path.join(ASSETS, 'astrometry_fig8_crop.png'),
                  5.6, 1.55, w=7.2, border=True)
caption(s, 'All 790 ECDFS tiles. Dashed: raw classical centroids. Solid: '
        'head-corrected — the clouds collapse and re-centre. Right: median '
        'offset against S/N.', 5.6, 1.55 + h + 0.14, 7.2)

s = blank_slide(prs, notes=(
    'And then better data arrived. Two independently Gaia-anchored '
    'solutions disagree by a coherent 9–10 mas — every arrow points the '
    'same way. That is not scatter; the landscape itself had moved. '
    'Nobody made a mistake. Say it plainly — this is the title of the '
    'talk, measured.'))
kicker(s, JAISP_KICK)
tf = textbox(s, 0.73, 1.6, 9.4)
lead(tf, [('Then better data arrived.', {'bold': True}),
          (' Two Gaia-anchored solutions disagree by a coherent ', {}),
          ('9–10 mas', {'bold': True}),
          (' — every arrow the same way.', {})], size=19, first=True,
     after=14)
lead(tf, 'Not scatter. The landscape itself had moved.', size=19)

# ---------------------------------------------------------- 9 · PHILOSOPHY
PHIL_KICK = 'Philosophy · The view from here'


def phil_slide(phrase, notes, first=False):
    s = blank_slide(prs, notes=notes)
    if first:
        scenehead(s, 'Philosophy', 'The view from here')
    else:
        kicker(s, PHIL_KICK)
    tf = textbox(s, 1.55, 2.15, 10.6)
    para(tf, phrase, size=26, color=INK, after=0, line=1.2, first=True)
    return s


s = phil_slide(
    'Every station of the loop was a choice.',
    notes=(
        'The long ascent — highest camera of the talk, almost all sky, the '
        'summit a speck at the bottom with him still sitting on it. Land '
        'the claim: every station of the loop was a choice, and choosing '
        'with reasons is philosophy. This is why a philosophy symposium at '
        'IPAC is not a joke.\n\nScene note (12-minute cut): beats 1, 5, 6 '
        'and 7 — choice, the sting, the rulers, the endings.'),
    first=True)
tf = textbox(s, 1.55, 3.15, 9.8)
lead(tf, 'The data we keep, the ruler we trust, the moment we stop — none '
     'of it is given by nature. Choosing with reasons is philosophy, and '
     'science runs on it quietly.', size=16, first=True)

s = phil_slide(
    'A relative loss steers an absolute science.',
    notes=(
        'The loss is relative — negotiated in committee, encoded in '
        'proposals and review panels. Reward a different better and '
        'different missions fly, different papers count: science climbs a '
        'different mountain. The direction of the field is '
        'loss-dependent.'))
tf = textbox(s, 1.55, 3.15, 9.8)
lead(tf, [('Better', {'upright': False, 'italic': True}),
          (' is negotiated — in panels, proposals, citations. Reward a '
           'different better, and different missions fly, different papers '
           'count: the field climbs a different mountain.', {})],
     size=16, first=True)

s = phil_slide(
    'An AI will be a different kind of scientist.',
    notes=(
        'The AI scientist, read off the table’s rows: it restarts; its '
        'data re-shuffle; its loss is explicit and optimizable. It will '
        'converge, efficiently, on exactly what we asked for — which is '
        'the danger. A brain with a muddled objective fails gently; a '
        'model with a crisp wrong objective converges precisely. The '
        'audit is ours.'))
tf = textbox(s, 1.55, 3.15, 9.8)
lead(tf, 'It restarts; its data re-shuffle; its loss is explicit — it will '
     'converge, efficiently, on exactly what we asked for.', size=16,
     first=True, after=12)
lead(tf, 'A brain fails gently. A model fails precisely. The audit is '
     'ours.', size=16)

s = phil_slide(
    'One grammar, different nouns.',
    notes=(
        'One grammar: priors, likelihoods, information, surprise — a first '
        'impression and a cosmology conjugate the same verbs. Measuring a '
        'person is not a metaphor for science; it is the same inference '
        'with different nouns.'))
tf = textbox(s, 1.55, 3.15, 9.8)
lead(tf, 'Priors, likelihoods, information, surprise: a first impression '
     'and a cosmology conjugate the same verbs. Measuring a person is not '
     'a metaphor for science — it is the same inference.', size=16,
     first=True)

s = phil_slide(
    'The tenth coffee, at field scale.',
    notes=(
        'The sting — deliver it dry: after ten meetings we hold a '
        'serviceable model of a person and stop. Yet we observe the same '
        'galaxy a thousand times, and write the same paper nine hundred. '
        'When the posterior stops moving and we keep observing anyway, the '
        'loss is speaking — citations, careers, committees — not the '
        'data.'))
tf = textbox(s, 1.55, 3.15, 9.8)
lead(tf, 'Ten meetings, and we hold a serviceable model of a person. Yet '
     'we observe the same galaxy a thousand times — and write the same '
     'paper nine hundred.', size=16, first=True, after=12)
lead(tf, 'When the posterior stops moving and we keep observing, the loss '
     'is speaking, not the data.', size=16)

s = phil_slide(
    'You carry rulers too.',
    notes=(
        'The rulers you carry: rationalist on questions of consistency, '
        'existentialist where deduction runs out and you must own the '
        'choice, absurdist when neither resolves. Not inconsistency — a '
        'different loss for a different question, and the swing between '
        'them is what a healthy updating system looks like. (If asked '
        'about the modern/postmodern oscillation by name: metamodernism, '
        'Vermeulen & van den Akker — keep it for Q&A, not the slide.)'))
mini_table(
    s, 1.55, 3.2, 10.2,
    ['rationalist', 'existentialist', 'absurdist'],
    [['where the question is one of consistency'],
     ['where deduction runs out and the choice is yours to own'],
     ['where neither of them resolves']])
tf = textbox(s, 1.55, 5.3, 9.8)
lead(tf, 'Not inconsistency: a different loss for a different question — '
     'and the swing between them is what an updating system looks like.',
     size=16, first=True)

s = blank_slide(prs, notes=(
    'Two ways for a question to end. Die on the hilltop: the answer '
    'sufficed; the summit was always local to the question. Die on the '
    'slope: the answer keeps improving, the theory stays incomplete, the '
    'experiment outlives you. Weinberg’s Dreams of a Final Theory — say '
    '“half in longing, half in mourning” as your phrase about reading '
    'him, not as his words.'))
kicker(s, PHIL_KICK)
tf = textbox(s, 0.73, 1.3, 9.0)
para(tf, 'Two ways for a question to end.', size=26, color=INK, after=0,
     line=1.2, first=True)
tf = textbox(s, 7.9, 2.6, 4.6)
lead(tf, [('Die on the hilltop', {'bold': True}),
          (' — the answer sufficed; the summit was always local to the '
           'question.', {})], size=16, first=True)
tf = textbox(s, 0.73, 4.4, 4.9)
lead(tf, [('Die on the slope', {'bold': True}),
          (' — the answer keeps improving, the theory stays incomplete, '
           'and the experiment outlives you.', {})], size=16, first=True,
     after=12)
lead(tf, [('Weinberg’s ', {}),
          ('Dreams of a Final Theory', {'italic': True}),
          (', read half in longing and half in mourning.', {})], size=16)

# -------------------------------------------------------------- 10 · RETURN
s = blank_slide(prs, notes=(
    'Exactly the opening camera — same hill, same scale, and by now it '
    'means something else. Sisyphus on his break, playing pinball: the '
    'ball comes back, and he paid for that. Let the drawing sit in '
    'silence for a moment.'))
picture(s, closing_png, 6.35, 1.15, h=5.2)

s = blank_slide(prs, notes='The last line. Then stop talking.')
picture(s, closing_png, 6.35, 1.15, h=5.2)
tf = textbox(s, 0.73, 3.1, 5.4)
para(tf, 'One must imagine Sisyphus happy.', size=28, color=INK, line=1.18,
     after=0, first=True)

prs.save(OUT)
print(f'wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides')
