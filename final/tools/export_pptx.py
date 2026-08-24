#!/usr/bin/env python3
"""
export_pptx.py — regenerate Presentation.pptx from the deck, with REAL text.

Per beat: captures an art-only background (headless Chrome, &notext=1) and
rebuilds every visible text block from scenes.js as a native, editable
PowerPoint text box at the same screen position, size and style. Speaker
notes carry the presenter notes.

Run:  final/tools/venv/bin/python final/tools/export_pptx.py
Output: Presentation.pptx at the repo root (overwritten; close it first).
"""

import json
import re
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path

REPO = Path(__file__).resolve().parents[2]
FINAL = REPO / "final"
OUT = REPO / "Presentation.pptx"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
SIZE = "3200,1800"

# stage px -> slide: 1600 px == 13.333 in  ->  1 px == 0.008333 in == 0.6 pt
IN_PER_PX = 13.333 / 1600.0
PT_PER_PX = 0.6

FONT = "Palatino"
INK = "1A2026"
INK_SOFT = "37424C"
MUTED = "66727C"
HUMAN = "A05A3C"

# ---- 1 · enumerate beats, notes, and per-beat text geometry from the deck.
# Mirrors the engine's camera/placement math (resolveCam, camFor, fx/fy).
ENUM_JS = r"""
global.window = {};
require(__TERRAIN__);
require(__SCENES_PATH__);
const S = window.SIS;
const resolve = (c) => {
  c = typeof c === 'function' ? c() : c;
  const o = {};
  for (const k of ['x','y','z']) o[k] = typeof c[k]==='function'?c[k]():c[k];
  return o;
};
const out = [];
S.SCENES.forEach((sc) => {
  const cam0 = resolve(sc.camera);
  const steps = sc.steps && sc.steps.length ? sc.steps : [{}];
  const cams = [];
  let cur = cam0;
  steps.forEach((st) => { if (st.camera) cur = resolve(st.camera); cams.push(cur); });
  steps.forEach((st, k) => {
    const texts = (sc.text || []).filter((t) => {
      const from = t.from || 0, to = t.to != null ? t.to : 99;
      return k >= from && k <= to;
    }).map((t) => {
      const cf = cams[Math.min(t.from || 0, steps.length - 1)];
      const ck = cams[k];
      const wx = cf.x - 800 / cf.z + (t.at[0] * 1600) / cf.z;
      const wy = cf.y - 450 / cf.z + (t.at[1] * 900) / cf.z;
      return {
        x: (wx - ck.x) * ck.z / 1600 + 0.5,
        y: (wy - ck.y) * ck.z / 900 + 0.5,
        scale: ck.z / cf.z,
        w: t.w, cls: t.cls || '', html: t.html,
      };
    });
    out.push({
      id: sc.id, name: sc.name, step: k, nsteps: steps.length,
      notes: (st.notes || sc.notes || '').trim(), texts,
    });
  });
});
console.log(JSON.stringify(out));
"""
ENUM_JS = ENUM_JS.replace("__TERRAIN__", repr(str(FINAL / "js" / "terrain.js")))
ENUM_JS = ENUM_JS.replace("__SCENES_PATH__", repr(str(FINAL / "js" / "scenes.js")))
beats = json.loads(subprocess.check_output(["node", "-e", ENUM_JS], text=True))
print(f"{len(beats)} beats across {len({b['id'] for b in beats})} scenes")

# ---- 2 · text styling ------------------------------------------------------
# paragraph styles keyed by (block class, inner tag/class)
def para_style(block_cls, inner):
    c = block_cls.split()
    s = {"size": 21, "italic": True, "color": MUTED, "bold": False,
         "caps": False, "spc": 0, "line": 1.4}
    if "titleblock" in c or "scenehead" in c:
        s.update(italic=False, color=INK)
        if inner == "kicker":
            s.update(size=15 if "titleblock" in c else 15, color=MUTED,
                     caps=True, spc=3, line=1.2)
        elif inner == "h1":
            s.update(size=72, line=1.04)
        elif inner == "scenetitle":
            s.update(size=52, line=1.06)
        elif inner == "byline":
            s.update(size=20, italic=True, color=INK_SOFT)
    elif "phrase" in c:
        s.update(size=36 if "sm" in c else 46, italic=False, color=INK, line=1.2)
    elif "mini-table" in c:
        if inner == "h5":
            s.update(size=13, italic=False, color=MUTED, caps=True, spc=2, line=1.3)
        else:
            s.update(size=19, italic=True, color=INK_SOFT, line=1.4)
    elif "bigtable" in c:
        if inner in ("h5", "h6-head"):
            s.update(size=12, italic=False, color=MUTED, caps=True, spc=2, line=1.3)
        elif inner == "h6":
            s.update(size=15.5, italic=False, color=INK, line=1.32)
        else:
            s.update(size=14.5, italic=True, color=INK, line=1.32)
    elif "aside" in c:
        if "lead" in c:
            s.update(size=26, color=INK, line=1.44)
        if "human" in c:
            s.update(color=HUMAN)
    return s

# ---- 3 · a tiny HTML-to-paragraphs parser ----------------------------------
class Mini(HTMLParser):
    """Flattens the deck's small HTML vocabulary into [(style, [runs])]."""
    def __init__(self, block_cls):
        super().__init__(convert_charrefs=True)
        self.block_cls = block_cls
        self.paras = []           # (inner_style_key, [ {text,bold,em,small,muted} ])
        self.cur = None
        self.bold = 0
        self.small = 0
    def handle_starttag(self, tag, attrs):
        a = dict(attrs)
        if tag in ("p", "h1", "h5", "h6"):
            inner = tag if tag != "p" else (a.get("class") or "p")
            self.cur = (inner, [])
            self.paras.append(self.cur)
        elif tag == "b":
            self.bold += 1
        elif tag in ("span", "sub"):
            self.small += 1
        elif tag == "br" and self.cur:
            self.cur = (self.cur[0], [])
            self.paras.append(self.cur)
    def handle_endtag(self, tag):
        if tag == "b":
            self.bold = max(0, self.bold - 1)
        elif tag in ("span", "sub"):
            self.small = max(0, self.small - 1)
    def handle_data(self, data):
        if not data or self.cur is None:
            return
        self.cur[1].append({
            "text": data.replace(" ", " "),
            "bold": self.bold > 0, "small": self.small > 0,
        })

def parse_block(cls, html):
    """Returns a list of columns; each column is a list of (style,[runs]).
    Grid blocks (bigtable / mini-table) split into columns; others one col."""
    cols = []
    if "bigtable" in cls or "mini-table" in cls:
        if "mini-table" in cls:
            parts = re.findall(r"<div>(.*?)</div>", html, re.S)
        else:
            # split the flat cell sequence: h6/h5 starts a new column only in
            # the header; rows are h6 + 3 <p> cells
            if "head" in cls:
                parts = re.findall(r"<h6>.*?</h6>", html, re.S)
                parts = [p.replace("h6>", "h6-head>") if False else p for p in parts]
            else:
                m = re.match(r"\s*(<h6>.*?</h6>)(.*)$", html, re.S)
                parts = [m.group(1)] + re.findall(r"<p>.*?</p>", m.group(2), re.S)
        for part in parts:
            p = Mini(cls)
            p.feed(part)
            if "bigtable" in cls and "head" in cls:
                p.paras = [("h6-head", runs) for _, runs in p.paras]
            cols.append(p.paras)
    else:
        p = Mini(cls)
        p.feed(html)
        cols.append(p.paras)
    return cols

# grid geometry (stage px, at scale 1) — must match deck.css
def grid_offsets(cls, w):
    if "mini-table" in cls:
        cw = (w - 2 * 34) / 3.0
        return [(i * (cw + 34), cw) for i in range(3)]
    if "bigtable" in cls:
        fr = (w - 150 - 3 * 22) / 3.0
        offs = [(0, 150)]
        x = 150 + 22
        for _ in range(3):
            offs.append((x, fr))
            x += fr + 22
        return offs
    return [(0, w)]

# ---- 4 · capture art-only backgrounds --------------------------------------
tmp = Path(tempfile.mkdtemp(prefix="deck-export-"))
shots = []
for i, b in enumerate(beats, 1):
    png = tmp / f"{i:02d}.png"
    url = (f"file://{FINAL}/index.html?scene={b['id']}&step={b['step']}"
           f"&instant=1&bare=1&notext=1")
    subprocess.run([CHROME, "--headless=new", "--disable-gpu",
                    f"--window-size={SIZE}", f"--screenshot={png}", url],
                   check=True, capture_output=True)
    jpg = png.with_suffix(".jpg")
    subprocess.run(["sips", "-s", "format", "jpeg", "-s", "formatOptions", "90",
                    str(png), "--out", str(jpg)], check=True, capture_output=True)
    shots.append(jpg)
    print(f"  [{i:2d}/{len(beats)}] {b['id']} · {b['step']+1}/{b['nsteps']}")

# ---- 5 · assemble -----------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_textbox(slide, x_in, y_in, w_in, paras, cls, scale):
    paras = [p for p in paras if p[1]]
    if not paras:
        return
    box = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                   Inches(w_in), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for inner, runs in paras:
        st = para_style(cls, inner)
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.line_spacing = st["line"]
        for r in runs:
            run = para.add_run()
            run.text = r["text"].upper() if st["caps"] else r["text"]
            f = run.font
            f.name = FONT
            px = st["size"] * scale                 # on-screen stage px
            if r["small"]:
                px = min(px, 17 * scale)
            f.size = Pt(max(8, px * PT_PER_PX))
            f.bold = r["bold"] or st["bold"]
            f.italic = st["italic"] and not r["bold"]
            f.color.rgb = RGBColor.from_string(MUTED if r["small"] else st["color"])
            if st["spc"]:
                f._rPr.set("spc", str(int(st["spc"] * 100)))

for jpg, b in zip(shots, beats):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(jpg), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    for t in b["texts"]:
        cols = parse_block(t["cls"], t["html"])
        offsets = grid_offsets(t["cls"], t["w"])
        for (off, cw), paras in zip(offsets, cols):
            x_in = (t["x"] * 1600 + off * t["scale"]) * IN_PER_PX
            y_in = t["y"] * 7.5
            add_textbox(slide, x_in, y_in, cw * t["scale"] * IN_PER_PX,
                        paras, t["cls"], t["scale"])
    head = f"{b['name']} — beat {b['step']+1}/{b['nsteps']}"
    slide.notes_slide.notes_text_frame.text = f"{head}\n\n{b['notes']}"

prs.save(OUT)
print(f"wrote {OUT.name}: {len(beats)} slides, {OUT.stat().st_size/1e6:.0f} MB")
